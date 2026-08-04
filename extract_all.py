import pptx
prs = pptx.Presentation('../TOEIC LISTENING - PART 01.pptx')
for idx in range(206, 212):
    slide = prs.slides[idx]
    for s_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "image"):
            with open(f"/tmp/ppt_debug/slide_{idx}_shape_{s_idx}.{shape.image.ext}", "wb") as f:
                f.write(shape.image.blob)
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for g_idx, gs in enumerate(shape.shapes):
                if hasattr(gs, "image"):
                    with open(f"/tmp/ppt_debug/slide_{idx}_group_{s_idx}_{g_idx}.{gs.image.ext}", "wb") as f:
                        f.write(gs.image.blob)
