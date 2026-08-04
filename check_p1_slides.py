import pptx
prs = pptx.Presentation('../TOEIC LISTENING - PART 01.pptx')
for idx in range(200, 250):
    slide = prs.slides[idx]
    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text += shape.text_frame.text + " "
    print(f"Slide {idx+1}: {text[:100].strip()}")
