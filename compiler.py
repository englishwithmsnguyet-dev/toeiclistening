import sys
import os
import re
import json
import zipfile
import html
from pptx import Presentation

# 1. Paths configuration
script_dir = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.abspath(os.path.join(script_dir, "..", "TOECI LISTENING - PART 03.pptx"))
output_dir = script_dir
data_dir = os.path.join(output_dir, "data")
media_dir = os.path.join(output_dir, "media")

# Create folders if not exist
os.makedirs(data_dir, exist_ok=True)
os.makedirs(media_dir, exist_ok=True)

print("Starting TOEIC Listening Part 3 HTML Compiler...")

# 2. Extract media from PPTX zip structure
with zipfile.ZipFile(pptx_path, 'r') as z:
    audio_files = [f for f in z.namelist() if 'media' in f and f.lower().endswith('.mp3')]
    print(f"Found {len(audio_files)} audio files in PPTX.")
    extracted_count = 0
    for f in audio_files:
        filename = os.path.basename(f)
        dest_path = os.path.join(media_dir, filename)
        if not os.path.exists(dest_path):
            with open(dest_path, 'wb') as dest:
                dest.write(z.read(f))
            extracted_count += 1
    print(f"Extracted {extracted_count} audio files to {media_dir}.")

# Load presentation
prs = Presentation(pptx_path)
slides = list(prs.slides)
print(f"Loaded presentation with {len(slides)} slides.")

# Helpers
def get_clean_lines(text_frame_text):
    lines = re.split(r'[\n\r\x0b]', text_frame_text)
    cleaned = []
    for l in lines:
        l_str = l.strip()
        if l_str:
            cleaned.append(l_str)
    return cleaned

def get_paragraph_html(paragraph):
    parts = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        
        # Escape HTML special characters
        t_escaped = html.escape(text)
        
        # Check run colors
        color = run.font.color
        has_color = False
        hex_color = ""
        if color and color.type == 1: # RGBColor
            try:
                hex_color = f"#{color.rgb}"
                has_color = True
            except:
                pass
                
        span_text = t_escaped
        style_attrs = []
        if has_color:
            style_attrs.append(f"color: {hex_color};")
            
        if style_attrs:
            style_str = " ".join(style_attrs)
            span_text = f'<span style="{style_str}">{span_text}</span>'
            
        if run.font.bold:
            span_text = f'<strong>{span_text}</strong>'
        if run.font.italic:
            span_text = f'<em>{span_text}</em>'
            
        # Replace vertical tab/soft returns with HTML breaks
        span_text = span_text.replace('\x0b', '<br>')
        span_text = span_text.replace('\u000b', '<br>')
        
        parts.append(span_text)
        
    return "".join(parts).strip()

def extract_slide_text(slide):
    # Sort shapes by their top coordinate to ensure top-to-bottom reading order
    sorted_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: s.top
    )
    texts = []
    for shape in sorted_shapes:
        for p in shape.text_frame.paragraphs:
            p_html = get_paragraph_html(p)
            if p_html:
                is_bullet = False
                pPr = p._p.pPr
                if pPr is not None:
                    has_bu_char = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar') is not None
                    has_bu_font = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buFont') is not None
                    has_bu_none = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone') is not None
                    if (has_bu_char or has_bu_font) and not has_bu_none:
                        is_bullet = True
                        
                # Fallback checking for manual bullets
                text_plain = p.text.strip()
                if not is_bullet:
                    if text_plain.startswith(('o ', '• ', '- ', '* ', '◦ ')):
                        is_bullet = True
                
                if is_bullet:
                    # Prepend standard bullet if it's a list item and doesn't already have one
                    if not text_plain.startswith(('o ', '• ', '- ', '* ', '◦ ')):
                        p_html = "• " + p_html
                        
                texts.append(p_html)
    return texts

def extract_slide_audio(slide):
    for shape in slide.shapes:
        xml_str = shape.element.xml
        if 'media' in xml_str or 'audio' in xml_str:
            rids = re.findall(r'r:(?:embed|id|link)="([^"]+)"', xml_str)
            for rid in rids:
                try:
                    target = slide.part.rels[rid].target_ref
                    if target.endswith('.mp3'):
                        return target.split('/')[-1]
                except:
                    pass
    return None

def find_nearest_audio(slides, target_slide_num):
    start = max(1, target_slide_num - 5)
    end = min(len(slides), target_slide_num + 5)
    for s_idx in range(start, end + 1):
        audio = extract_slide_audio(slides[s_idx - 1])
        if audio:
            return audio
    return None

def parse_question_slide(slide, slide_num):
    tb = None
    for shape in slide.shapes:
        if shape.has_text_frame and "A." in shape.text_frame.text:
            tb = shape
            break
            
    if not tb:
        return None
        
    p_htmls = []
    for p in tb.text_frame.paragraphs:
        p_html = get_paragraph_html(p)
        if p_html:
            p_htmls.append(p_html)
            
    # Split visual lines by breaks and track column index
    v_lines_with_col = []
    for p_html in p_htmls:
        lines = re.split(r'<br\s*/?>', p_html)
        for line in lines:
            subparts = re.split(r'\t|\s{2,}(?=[A-D]\.)|\s+(?=[B-D]\.)', line)
            subparts = [s.strip() for s in subparts if s.strip()]
            for col_idx, part in enumerate(subparts):
                v_lines_with_col.append((part, col_idx))
                
    question_text = ""
    choices = {}
    choice_lines = []
    
    for l_idx, (line_html, col_idx) in enumerate(v_lines_with_col):
        # Match choice letters robustly even if wrapped in HTML tags
        text_plain = re.sub(r'<[^>]+>', '', line_html).strip()
        match = re.match(r'^([A-D])\.\s*(.*)', text_plain)
        if match:
            letter = match.group(1)
            # Remove "A. " indicator but preserve preceding HTML tags if any
            html_clean = re.sub(r'^((?:<[^>]+>)*)([A-D])\.\s*', r'\1', line_html)
            choices[letter] = html_clean
            choice_lines.append((l_idx, letter, col_idx))
        else:
            if not question_text:
                question_text = line_html
            else:
                question_text += "<br>" + line_html
                
    # Resolve answers using Oval shape
    oval_shape = None
    for shape in slide.shapes:
        if "oval" in shape.name.lower() or "circle" in shape.name.lower():
            oval_shape = shape
            break
            
    correct_answer = None
    if oval_shape and choice_lines:
        tb_top = tb.top
        tb_height = tb.height
        tb_left = tb.left
        tb_width = tb.width
        oval_top = oval_shape.top
        oval_left = oval_shape.left
        
        # Calculate centers
        oval_x = oval_left + oval_shape.width / 2
        
        # Check if oval is horizontally in the left half or right half of the text box
        is_left_side = oval_x < (tb_left + tb_width * 0.48)
        target_col = 0 if is_left_side else 1
        
        # Check if we have choices in the target column
        has_target_col_choices = any(col == target_col for _, _, col in choice_lines)
        
        best_letter = None
        min_dist = float('inf')
        for l_idx, letter, col_idx in choice_lines:
            # If layout has choices in the correct column, filter candidates
            if has_target_col_choices:
                if col_idx != target_col:
                    continue
                    
            y_est = tb_top + ((l_idx + 0.5) / len(v_lines_with_col)) * tb_height
            dist = abs(oval_top - y_est)
            if dist < min_dist:
                min_dist = dist
                best_letter = letter
        correct_answer = best_letter
        
    audio = extract_slide_audio(slide)
    
    return {
        "slide_index": slide_num,
        "question": question_text,
        "choices": choices,
        "answer": correct_answer,
        "audio": audio
    }

def get_slide_script_html(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if any(text.startswith(marker) for marker in ["W:", "M:", "W-", "M-", "W1:", "W2:", "M1:", "M2:"]):
                p_htmls = []
                for p in shape.text_frame.paragraphs:
                    p_html = get_paragraph_html(p)
                    if p_html:
                        p_htmls.append(p_html)
                return p_htmls
    return []

# Structure definition (1-based slide indices)
# Classify theory and vocabulary slides separately
structure_map = [
    # --- OVERVIEW ---
    {
        "id": "overview",
        "title": "I. TỔNG QUAN PHẦN 03",
        "type": "overview",
        "slides": [2, 3]
    },
    
    # --- DẠNG 01: CÂU HỎI TỔNG QUAN ---
    {
        "id": "dang_01_topic",
        "title": "Dạng 1.I - Câu hỏi về Chủ đề bài nghe",
        "type": "subsection",
        "parent_title": "DẠNG 01: CÂU HỎI TỔNG QUAN",
        "theory_slides": [4, 5],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 6, "s": 7},
            {"title": "Ví dụ 2", "q": 8, "s": 9}
        ],
        "practice_slides": [
            {"q": 11, "s": 12}, {"q": 13, "s": 14}, {"q": 15, "s": 16}, {"q": 17, "s": 18}, {"q": 19, "s": 20},
            {"q": 21, "s": 22}, {"q": 23, "s": 24}, {"q": 25, "s": 26}, {"q": 27, "s": 28}, {"q": 29, "s": 30}
        ]
    },
    {
        "id": "dang_01_location",
        "title": "Dạng 1.II - Câu hỏi về Địa điểm",
        "type": "subsection",
        "parent_title": "DẠNG 01: CÂU HỎI TỔNG QUAN",
        "theory_slides": [31, 32],
        "vocabulary_slides": list(range(33, 47)), # 33 to 46
        "example_slides": [
            {"title": "Ví dụ 1", "q": 47, "s": 48},
            {"title": "Ví dụ 2", "q": 49, "s": 50}
        ],
        "practice_slides": [
            {"q": 52, "s": 53}, {"q": 54, "s": 55}, {"q": 56, "s": 57}, {"q": 58, "s": 59}, {"q": 60, "s": 61},
            {"q": 62, "s": 63}, {"q": 64, "s": 65}, {"q": 66, "s": 67}, {"q": 68, "s": 69}, {"q": 70, "s": 71}
        ]
    },
    {
        "id": "dang_01_identity",
        "title": "Dạng 1.III - Câu hỏi về Nghề nghiệp",
        "type": "subsection",
        "parent_title": "DẠNG 01: CÂU HỎI TỔNG QUAN",
        "theory_slides": [72, 73],
        "vocabulary_slides": list(range(74, 77)), # 74 to 76
        "example_slides": [
            {"title": "Ví dụ 1", "q": 77, "s": 78},
            {"title": "Ví dụ 2", "q": 79, "s": 80},
            {"title": "Ví dụ 3", "q": 81, "s": 82}
        ],
        "practice_slides": [
            {"q": 84, "s": 85}, {"q": 86, "s": 87}, {"q": 88, "s": 89}, {"q": 90, "s": 91}, {"q": 92, "s": 93},
            {"q": 94, "s": 95}, {"q": 96, "s": 97}, {"q": 98, "s": 99}, {"q": 100, "s": 101}, {"q": 102, "s": 103}
        ]
    },
    
    # --- DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT ---
    {
        "id": "dang_02_problem",
        "title": "Dạng 2.I - Câu hỏi về Vấn đề / Nguyên nhân",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT",
        "theory_slides": [104, 105],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 106, "s": 107},
            {"title": "Ví dụ 2", "q": 108, "s": 109},
            {"title": "Ví dụ 3", "q": 110, "s": 111}
        ],
        "practice_slides": [
            {"q": 113, "s": 114}, {"q": 115, "s": 116}, {"q": 117, "s": 118}, {"q": 119, "s": 120}, {"q": 121, "s": 122},
            {"q": 123, "s": 124}, {"q": 125, "s": 126}, {"q": 127, "s": 128}, {"q": 129, "s": 130}, {"q": 131, "s": 132}
        ]
    },
    {
        "id": "dang_02_why",
        "title": "Dạng 2.II - Câu hỏi Why",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT",
        "theory_slides": [133, 134, 135],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 136, "s": 137},
            {"title": "Ví dụ 2", "q": 138, "s": 139}
        ],
        "practice_slides": [
            {"q": 141, "s": 142}, {"q": 143, "s": 144}, {"q": 145, "s": 146}, {"q": 147, "s": 148}, {"q": 149, "s": 150},
            {"q": 151, "s": 152}, {"q": 153, "s": 154}, {"q": 155, "s": 156}, {"q": 157, "s": 158}, {"q": 159, "s": 160}
        ]
    },
    {
        "id": "dang_02_what_request",
        "title": "Dạng 2.III.1 - Câu hỏi Yêu cầu",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT - CÂU HỎI WHAT",
        "theory_slides": [161, 162],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 163, "s": 164},
            {"title": "Ví dụ 2", "q": 165, "s": 166}
        ],
        "practice_slides": [
            {"q": 168, "s": 169}, {"q": 170, "s": 171}, {"q": 172, "s": 173}
        ]
    },
    {
        "id": "dang_02_what_suggest",
        "title": "Dạng 2.III.2 - Câu hỏi Gợi ý",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT - CÂU HỎI WHAT",
        "theory_slides": [174, 175],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 176, "s": 177},
            {"title": "Ví dụ 2", "q": 178, "s": 179}
        ],
        "practice_slides": [
            {"q": 181, "s": 182}, {"q": 183, "s": 184}, {"q": 185, "s": 186}
        ]
    },
    {
        "id": "dang_02_what_say",
        "title": "Dạng 2.III.3 - Câu hỏi 'What does the speaker say?'",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT - CÂU HỎI WHAT",
        "theory_slides": [187],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 188, "s": 189},
            {"title": "Ví dụ 2", "q": 190, "s": 191}
        ],
        "practice_slides": [
            {"q": 193, "s": 194}, {"q": 195, "s": 196}, {"q": 197, "s": 198}
        ]
    },
    {
        "id": "dang_02_what_according",
        "title": "Dạng 2.III.4 - Câu hỏi 'According to...'",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT - CÂU HỎI WHAT",
        "theory_slides": [199, 200],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 201, "s": 202},
            {"title": "Ví dụ 2", "q": 203, "s": 204}
        ],
        "practice_slides": [
            {"q": 206, "s": 207}, {"q": 208, "s": 209}, {"q": 210, "s": 211}
        ]
    },
    {
        "id": "dang_02_what_do_next",
        "title": "Dạng 2.III.5 - Câu hỏi 'Do next...'",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT - CÂU HỎI WHAT",
        "theory_slides": [212, 213],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 214, "s": 215},
            {"title": "Ví dụ 2", "q": 216, "s": 217}
        ],
        "practice_slides": [
            {"q": 219, "s": 220}, {"q": 221, "s": 222}, {"q": 223, "s": 224}
        ]
    },
    {
        "id": "dang_02_what_imply",
        "title": "Dạng 2.III.6 - Câu hỏi Ngụ ý",
        "type": "subsection",
        "parent_title": "DẠNG 02: CÂU HỎI VỀ THÔNG TIN CHI TIẾT - CÂU HỎI WHAT",
        "theory_slides": [225, 226],
        "vocabulary_slides": [],
        "example_slides": [
            {"title": "Ví dụ 1", "q": 227, "s": 228},
            {"title": "Ví dụ 2", "q": 229, "s": 230}
        ],
        "practice_slides": [
            {"q": 232, "s": 233}, {"q": 234, "s": 235}, {"q": 236, "s": 237}
        ]
    },
    {
        "id": "dang_visual_questions",
        "title": "V. KẾT HỢP HÌNH ẢNH - Luyện tập",
        "type": "subsection",
        "parent_title": "CÂU HỎI KẾT HỢP HÌNH ẢNH",
        "theory_slides": [],
        "vocabulary_slides": [],
        "example_slides": [],
        "practice_slides": [
            {"q": 327, "s": 328},
            {"q": 336, "s": 337},
            {"q": 436, "s": 437},
            {"q": 441, "s": 442},
            {"q": 455, "s": 456},
            {"q": 481, "s": 482},
            {"q": 490, "s": 491},
            {"q": 540, "s": 541}
        ]
    },
    
    # --- LƯU Ý KHI LÀM BÀI ---
    {
        "id": "tips",
        "title": "Lưu ý khi làm bài Listening Part 03",
        "type": "tips",
        "slides": list(range(238, 245))
    },
    
    # --- CHỦ ĐỀ LUYỆN TẬP (3-QUESTION SETS) ---
    {
        "id": "topic_01",
        "title": "Chủ đề 1 - Nhân sự trong công ty",
        "type": "topic",
        "theory_slides": [246, 247, 248],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [249, 250, 251, 252]},
            {"slides": [253, 254, 255, 256]}
        ],
        "practice_sets": [
            {"slides": list(range(258, 265))},
            {"slides": list(range(265, 272))},
            {"slides": list(range(272, 279))},
            {"slides": list(range(279, 286))}
        ]
    },
    {
        "id": "topic_02",
        "title": "Chủ đề 2 - Cuộc họp, hội nghị, hội thảo",
        "type": "topic",
        "theory_slides": [287],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [288, 289, 290, 291]},
            {"slides": [292, 293, 294, 295]}
        ],
        "practice_sets": [
            {"slides": list(range(297, 304))},
            {"slides": list(range(304, 311))},
            {"slides": list(range(311, 318))},
            {"slides": list(range(318, 325))},
            {"slides": list(range(325, 332))},
            {"slides": list(range(332, 339))}
        ]
    },
    {
        "id": "topic_03",
        "title": "Chủ đề 3 - Những vấn đề kỹ thuật",
        "type": "topic",
        "theory_slides": [340, 341],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [342, 343, 344, 345]},
            {"slides": [346, 347, 348, 349]}
        ],
        "practice_sets": [
            {"slides": list(range(351, 358))},
            {"slides": list(range(358, 365))},
            {"slides": list(range(365, 372))},
            {"slides": list(range(372, 379))},
            {"slides": list(range(379, 386))},
            {"slides": list(range(386, 393))}
        ]
    },
    {
        "id": "topic_04",
        "title": "Chủ đề 4 - Đi lại, du lịch",
        "type": "topic",
        "theory_slides": [394],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [395, 396, 397, 398]},
            {"slides": [399, 400, 401, 402]}
        ],
        "practice_sets": [
            {"slides": list(range(404, 411))},
            {"slides": list(range(411, 418))},
            {"slides": list(range(418, 425))},
            {"slides": list(range(425, 432))},
            {"slides": list(range(432, 439))},
            {"slides": list(range(439, 446))}
        ]
    },
    {
        "id": "topic_05",
        "title": "Chủ đề 5 - Mua sắm",
        "type": "topic",
        "theory_slides": [447, 448],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [449, 450, 451, 452]},
            {"slides": [453, 454, 455, 456]}
        ],
        "practice_sets": [
            {"slides": list(range(458, 465))},
            {"slides": list(range(465, 472))},
            {"slides": list(range(472, 479))},
            {"slides": list(range(479, 486))},
            {"slides": list(range(486, 493))},
            {"slides": list(range(493, 500))}
        ]
    },
    {
        "id": "topic_06",
        "title": "Chủ đề 6 - Giải trí",
        "type": "topic",
        "theory_slides": [501],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [502, 503, 504, 505]},
            {"slides": [506, 507, 508, 509]}
        ],
        "practice_sets": [
            {"slides": list(range(511, 518))},
            {"slides": list(range(518, 525))},
            {"slides": list(range(525, 532))},
            {"slides": [532, 533, 534, 535, 536, 536, 537]},
            {"slides": list(range(538, 545))},
            {"slides": list(range(545, 552))}
        ]
    }
]

# 4. Processing Slide Content
final_data = []

for entry in structure_map:
    print(f"Processing Section: {entry['title']}...")
    
    if entry["type"] == "overview" or entry["type"] == "tips":
        theory_data = []
        for s_num in entry["slides"]:
            slide = slides[s_num - 1]
            theory_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
        final_data.append({
            "id": entry["id"],
            "title": entry["title"],
            "type": entry["type"],
            "theory": theory_data,
            "vocabulary": []
        })
        
    elif entry["type"] == "subsection":
        theory_data = []
        for s_num in entry["theory_slides"]:
            slide = slides[s_num - 1]
            theory_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
            
        vocabulary_data = []
        for s_num in entry["vocabulary_slides"]:
            slide = slides[s_num - 1]
            vocabulary_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
            
        examples_data = []
        for ex in entry["example_slides"]:
            q_slide = slides[ex["q"] - 1]
            s_slide = slides[ex["s"] - 1]
            
            q_info = parse_question_slide(q_slide, ex["q"])
            if q_info:
                if not q_info["audio"]:
                    q_info["audio"] = extract_slide_audio(s_slide)
                if not q_info["audio"]:
                    q_info["audio"] = find_nearest_audio(slides, ex["q"])
                q_info["title"] = ex["title"]
                q_info["transcript"] = get_slide_script_html(s_slide)
                examples_data.append(q_info)
                
        practice_data = []
        for pr_pair in entry["practice_slides"]:
            q_slide = slides[pr_pair["q"] - 1]
            s_slide = slides[pr_pair["s"] - 1]
            
            q_info = parse_question_slide(q_slide, pr_pair["q"])
            if q_info:
                if not q_info["audio"]:
                    q_info["audio"] = extract_slide_audio(s_slide)
                if not q_info["audio"]:
                    q_info["audio"] = find_nearest_audio(slides, pr_pair["q"])
                q_info["transcript"] = get_slide_script_html(s_slide)
                practice_data.append(q_info)
                
        final_data.append({
            "id": entry["id"],
            "title": entry["title"],
            "parent_title": entry["parent_title"],
            "type": "subsection",
            "theory": theory_data,
            "vocabulary": vocabulary_data,
            "examples": examples_data,
            "practice": practice_data
        })
        
    elif entry["type"] == "topic":
        theory_data = []
        for s_num in entry["theory_slides"]:
            slide = slides[s_num - 1]
            theory_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
            
        vocabulary_data = []
        for s_num in entry["vocabulary_slides"]:
            slide = slides[s_num - 1]
            vocabulary_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
            
        example_sets_data = []
        for set_idx, ex_set in enumerate(entry.get("example_sets", [])):
            set_slides_nums = ex_set["slides"]
            
            q1_num = set_slides_nums[0]
            q2_num = set_slides_nums[1]
            q3_num = set_slides_nums[2]
            full_s_num = set_slides_nums[3]
            
            q1_info = parse_question_slide(slides[q1_num - 1], q1_num)
            q2_info = parse_question_slide(slides[q2_num - 1], q2_num)
            q3_info = parse_question_slide(slides[q3_num - 1], q3_num)
            
            full_transcript = extract_slide_text(slides[full_s_num - 1])
            
            # Scan all slides in the set to find the audio file
            audio = None
            for s_num in set_slides_nums:
                slide_audio = extract_slide_audio(slides[s_num - 1])
                if slide_audio:
                    audio = slide_audio
                    break
                    
            questions = []
            for i, qi in enumerate([q1_info, q2_info, q3_info]):
                if qi:
                    questions.append({
                        "id": i + 1,
                        "slide_index": qi["slide_index"],
                        "question": qi["question"],
                        "choices": qi["choices"],
                        "answer": qi["answer"]
                    })
                    
            example_sets_data.append({
                "set_index": set_idx + 1,
                "audio": audio,
                "questions": questions,
                "transcript": full_transcript
            })

        practice_sets_data = []
        for set_idx, pr_set in enumerate(entry["practice_sets"]):
            set_slides_nums = pr_set["slides"]
            
            q1_num = set_slides_nums[0]
            s1_num = set_slides_nums[1]
            q2_num = set_slides_nums[2]
            s2_num = set_slides_nums[3]
            q3_num = set_slides_nums[4]
            s3_num = set_slides_nums[5]
            full_s_num = set_slides_nums[6]
            
            q1_info = parse_question_slide(slides[q1_num - 1], q1_num)
            q2_info = parse_question_slide(slides[q2_num - 1], q2_num)
            q3_info = parse_question_slide(slides[q3_num - 1], q3_num)
            
            full_transcript = extract_slide_text(slides[full_s_num - 1])
            
            # Scan all slides in the set to find the audio file
            audio = None
            for s_num in set_slides_nums:
                slide_audio = extract_slide_audio(slides[s_num - 1])
                if slide_audio:
                    audio = slide_audio
                    break
                
            questions = []
            for i, qi in enumerate([q1_info, q2_info, q3_info]):
                if qi:
                    questions.append({
                        "id": i + 1,
                        "slide_index": qi["slide_index"],
                        "question": qi["question"],
                        "choices": qi["choices"],
                        "answer": qi["answer"]
                    })
                    
            practice_sets_data.append({
                "set_index": set_idx + 1,
                "audio": audio,
                "questions": questions,
                "transcript": full_transcript
            })
            
        final_data.append({
            "id": entry["id"],
            "title": entry["title"],
            "type": "topic",
            "theory": theory_data,
            "vocabulary": vocabulary_data,
            "examples": example_sets_data,
            "practice_sets": practice_sets_data
        })

# 5. Output file writing
output_json_path = os.path.join(data_dir, "part03_data.json")
raw_sets = [
    {
        "q_num": "32-34",
        "audio": "E26-T01-32-34.mp3",
        "transcript": [
            "W: Hey, Oliver. Did you see the focus group results for (32) our new spicy cheddar cheese? Everyone really liked it.",
            "M: Yes. It should be a great addition to our company's line of cheeses.",
            "W: Several people mentioned that they'd like to use it in recipes-to add to sauces, for example.",
            "M: (33) So maybe we should consider selling a shredded version that would melt easily when cooked.",
            "W: I'm sure we could do that. (34) I'll get in touch with the production manager with that request."
        ],
        "questions": [
            {
                "id": 32,
                "question": "<strong>PRACTICE 01</strong><br><strong>Question 32.</strong> What type of food product does the speakers' company sell?",
                "choices": {
                    "A": "Candy",
                    "B": "Cheese",
                    "C": "Bread",
                    "D": "Pasta"
                },
                "answer": "B"
            },
            {
                "id": 33,
                "question": "<strong>Question 33.</strong> What does the man suggest?",
                "choices": {
                    "A": "Lowering prices",
                    "B": "Hiring more workers",
                    "C": "Publishing a recipe",
                    "D": "Offering additional options"
                },
                "answer": "D"
            },
            {
                "id": 34,
                "question": "<strong>Question 34.</strong> What does the woman say she will do?",
                "choices": {
                    "A": "Send a schedule update",
                    "B": "Contact a production manager",
                    "C": "Visit the company headquarters",
                    "D": "Plan an advertising campaign"
                },
                "answer": "B"
            }
        ]
    },
    {
        "q_num": "35-37",
        "audio": "E26-T01-35-37.mp3",
        "transcript": [
            "M: Hi. (35) I'm calling to book three tickets for this Thursday's tennis match. Are there any seats left?",
            "W: Just a few! Tickets for Thursday's match have been selling quickly.",
            "M: I'm not surprised! After all, (36) Ife Rotimi won the regional championship tournament last month. Everyone wants to see her play after her incredible performance. What seats are available?",
            "W: Well, there's only one group of three seats together. (37) Advance payment is required to hold them."
        ],
        "questions": [
            {
                "id": 35,
                "question": "<strong>PRACTICE 02</strong><br><strong>Question 35.</strong> Why is the man calling?",
                "choices": {
                    "A": "To sign up for lessons",
                    "B": "To enter a competition",
                    "C": "To buy tickets to an event",
                    "D": "To ask about branded merchandise"
                },
                "answer": "C"
            },
            {
                "id": 36,
                "question": "<strong>Question 36.</strong> What did Ife Rotimi do last month?",
                "choices": {
                    "A": "She won a regional tournament.",
                    "B": "She gave a television interview.",
                    "C": "She started an institute.",
                    "D": "She hired a new coach."
                },
                "answer": "A"
            },
            {
                "id": 37,
                "question": "<strong>Question 37.</strong> What does the woman say is required?",
                "choices": {
                    "A": "A parking permit",
                    "B": "A photo ID",
                    "C": "Contact information",
                    "D": "Advance payment"
                },
                "answer": "D"
            }
        ]
    },
    {
        "q_num": "38-40",
        "audio": "E26-T01-38-40.mp3",
        "transcript": [
            "W: (38) Thanks for agreeing to help me organize the library's annual fund-raising dinner, Klaus. We hope the event brings in enough money to expand our children's book section.",
            "M: What task would you like me to start with?",
            "W: (39) Well, I could use some help sending out the invitations.",
            "M: OK, I can take care of that. (40) Is there a list of attendees available?",
            "W: It's in my computer files. (40) I'll e-mail it to you."
        ],
        "questions": [
            {
                "id": 38,
                "question": "<strong>PRACTICE 03</strong><br><strong>Question 38.</strong> What event are the speakers planning?",
                "choices": {
                    "A": "A fund-raising dinner",
                    "B": "An art gallery opening",
                    "C": "An awards ceremony",
                    "D": "A children's book fair"
                },
                "answer": "A"
            },
            {
                "id": 39,
                "question": "<strong>Question 39.</strong> What task does the woman ask the man to help with?",
                "choices": {
                    "A": "Arranging a shuttle service",
                    "B": "Choosing a catering firm",
                    "C": "Preparing a speech",
                    "D": "Sending out invitations"
                },
                "answer": "D"
            },
            {
                "id": 40,
                "question": "<strong>Question 40.</strong> What does the woman say she will do?",
                "choices": {
                    "A": "E-mail a list",
                    "B": "Speak with a colleague",
                    "C": "Provide a password",
                    "D": "Post a job opening"
                },
                "answer": "A"
            }
        ]
    },
    {
        "q_num": "41-43",
        "audio": "E26-T01-41-43.mp3",
        "transcript": [
            "W: Hey, Brian and Matteo. (41) I found some great pens to give away at the community festival to promote our business.",
            "M1: Great. Can we put our cleaning service logo on them?",
            "W: Yes, for no extra charge. And they're biodegradable. (42) They're made from paper.",
            "M2: So when we hand them out, we can mention that.",
            "M1: As well as talk about the organic cleaning supplies our company uses.",
            "W: OK. (43) I'll go ahead and order several cases."
        ],
        "questions": [
            {
                "id": 41,
                "question": "<strong>PRACTICE 04</strong><br><strong>Question 41.</strong> What event are the speakers preparing for?",
                "choices": {
                    "A": "A new-employee orientation",
                    "B": "A grand opening",
                    "C": "A community festival",
                    "D": "A trade show"
                },
                "answer": "C"
            },
            {
                "id": 42,
                "question": "<strong>Question 42.</strong> What is mentioned about some pens?",
                "choices": {
                    "A": "They are available in multiple colors.",
                    "B": "They use permanent ink.",
                    "C": "They are preferred by book authors.",
                    "D": "They are made from paper."
                },
                "answer": "D"
            },
            {
                "id": 43,
                "question": "<strong>Question 43.</strong> What does the woman offer to do?",
                "choices": {
                    "A": "Reserve a booth",
                    "B": "Place an order",
                    "C": "Organize a focus group",
                    "D": "Revise a budget"
                },
                "answer": "B"
            }
        ]
    },
    {
        "q_num": "44-46",
        "audio": "E26-T01-44-46.mp3",
        "transcript": [
            "W: (44) Jamestown Recycling Facility. How can I help you?",
            "M: Hi. I'm preparing to move soon, and (45) I have some electronics, such as televisions and computers, that I'd like to get rid of before I put my house on the market. My friend mentioned you might take them.",
            "W: Yes, that's right. We'll take all electronics.",
            "M: Great. I just have one question. Do you provide a pickup service?",
            "W: No, unfortunately you'll have to bring everything here yourself. However, (46) on our Web site we list a number of companies that can remove and dispose of the items for you."
        ],
        "questions": [
            {
                "id": 44,
                "question": "<strong>PRACTICE 05</strong><br><strong>Question 44.</strong> Where does the woman work?",
                "choices": {
                    "A": "At a delivery service",
                    "B": "At an electronics store",
                    "C": "At a recycling facility",
                    "D": "At a real estate agency"
                },
                "answer": "C"
            },
            {
                "id": 45,
                "question": "<strong>Question 45.</strong> What does the man want to dispose of?",
                "choices": {
                    "A": "Yard waste",
                    "B": "Used furniture",
                    "C": "Electronics",
                    "D": "Books"
                },
                "answer": "C"
            },
            {
                "id": 46,
                "question": "<strong>Question 46.</strong> What does the woman say can be found on a Web site?",
                "choices": {
                    "A": "A list of companies",
                    "B": "Hours of operation",
                    "C": "A permit application",
                    "D": "Directions to a site"
                },
                "answer": "A"
            }
        ]
    },
    {
        "q_num": "47-49",
        "audio": "E26-T01-47-49.mp3",
        "transcript": [
            "M: Zaina! What a surprise! (47) I haven't seen you since we took that class for business owners together last year. How are you?",
            "W: Great, thanks. I was just in the neighborhood and (48) thought I'd stop in for a cookie or a piece of cake. You have so many delicious baked goods here.",
            "M: Thank you! It's been a good year for business. I'm even considering opening a second location.",
            "W: Really? Well, I noticed that Sunnyvale Restaurant went out of business, and the building's up for lease. (49) It's very close to the local university. You'd probably get a lot of walk-in customers."
        ],
        "questions": [
            {
                "id": 47,
                "question": "<strong>PRACTICE 06</strong><br><strong>Question 47.</strong> How do the speakers know each other?",
                "choices": {
                    "A": "They took a class together.",
                    "B": "They used to work for the same company.",
                    "C": "They grew up in the same neighborhood.",
                    "D": "They met on a train."
                },
                "answer": "A"
            },
            {
                "id": 48,
                "question": "<strong>Question 48.</strong> What type of business does the man most likely own?",
                "choices": {
                    "A": "A fitness center",
                    "B": "A real estate agency",
                    "C": "A culinary school",
                    "D": "A bakery"
                },
                "answer": "D"
            },
            {
                "id": 49,
                "question": "<strong>Question 49.</strong> What advantage does the woman point out about a rental space?",
                "choices": {
                    "A": "Its price",
                    "B": "Its size",
                    "C": "Its location",
                    "D": "Its design"
                },
                "answer": "C"
            }
        ]
    },
    {
        "q_num": "50-52",
        "audio": "E26-T01-50-52.mp3",
        "transcript": [
            "W: Hi, Koji. (50) I think our new video game is nearly ready to be released. Are you aware of any improvements that need to be made before then?",
            "M: Actually, (51) I just finished testing the game this morning. I found a problem in the third stage of the game. There were a few times when my character couldn't move.",
            "W: Oh, that's strange!",
            "M: I double-checked the problem using a different controller. The same issue came up.",
            "W: Oh. (52) I think Pauline had a similar problem with a game she tested. Maybe you should ask her about it."
        ],
        "questions": [
            {
                "id": 50,
                "question": "<strong>PRACTICE 07</strong><br><strong>Question 50.</strong> Who most likely are the speakers?",
                "choices": {
                    "A": "Film actors",
                    "B": "Museum directors",
                    "C": "Video game developers",
                    "D": "Investigative journalists"
                },
                "answer": "C"
            },
            {
                "id": 51,
                "question": "<strong>Question 51.</strong> What did the man recently do?",
                "choices": {
                    "A": "He secured some funding.",
                    "B": "He tested a product.",
                    "C": "He read a script.",
                    "D": "He conducted an interview."
                },
                "answer": "B"
            },
            {
                "id": 52,
                "question": "<strong>Question 52.</strong> What does the woman suggest?",
                "choices": {
                    "A": "Consulting a colleague",
                    "B": "Planning an event",
                    "C": "Negotiating a contract",
                    "D": "Giving a client an update"
                },
                "answer": "A"
            }
        ]
    },
    {
        "q_num": "53-55",
        "audio": "E26-T01-53-55.mp3",
        "transcript": [
            "M: (53) You've reached the maintenance office at Hillview Apartment Complex.",
            "W: Hi. This is Palavi Sen from unit 35B. (54) I'm calling because the new thermostat in my apartment isn't working. It keeps shutting off and turning on randomly, so my apartment is getting cold.",
            "M: When did this issue start?",
            "W: A few hours ago. The thermostat was just installed yesterday.",
            "M: OK. (55) I can come and take a look at it tomorrow morning.",
            "W: But it's supposed to be below freezing tonight!"
        ],
        "questions": [
            {
                "id": 53,
                "question": "<strong>PRACTICE 08</strong><br><strong>Question 53.</strong> Who most likely is the man?",
                "choices": {
                    "A": "A delivery driver",
                    "B": "A security guard",
                    "C": "A maintenance worker",
                    "D": "A customer service representative"
                },
                "answer": "C"
            },
            {
                "id": 54,
                "question": "<strong>Question 54.</strong> What problem does the woman describe?",
                "choices": {
                    "A": "A device is malfunctioning.",
                    "B": "A key is missing.",
                    "C": "A parking area is unavailable.",
                    "D": "A package was not received."
                },
                "answer": "A"
            },
            {
                "id": 55,
                "question": "<strong>Question 55.</strong> What does the woman mean when she says, \"it's supposed to be below freezing tonight\"?",
                "choices": {
                    "A": "She is surprised by the weather forecast.",
                    "B": "She wants a service to be completed sooner.",
                    "C": "She will move some items indoors.",
                    "D": "She would prefer to park near her apartment."
                },
                "answer": "B"
            }
        ]
    },
    {
        "q_num": "56-58",
        "audio": "E26-T01-56-58.mp3",
        "transcript": [
            "W: Good morning! Welcome to Jasper Bank.",
            "M1: (56) Thanks for meeting with us to discuss a loan for our business.",
            "W: Why don't you tell me more about your business? I understand it's a repair shop?",
            "M2: Well, ten years ago, we opened as a snowmobile repair shop, but after a few years, (57) we also started renting out snowmobiles and other sports equipment.",
            "M1: Yes, and (58) because winter tourism has increased recently, we'd like to expand our space so that we can carry more inventory."
        ],
        "questions": [
            {
                "id": 56,
                "question": "<strong>PRACTICE 09</strong><br><strong>Question 56.</strong> Why do the men want to speak to the woman?",
                "choices": {
                    "A": "To review a building design",
                    "B": "To discuss a loan",
                    "C": "To develop an advertising plan",
                    "D": "To purchase some supplies"
                },
                "answer": "B"
            },
            {
                "id": 57,
                "question": "<strong>Question 57.</strong> What type of business do the men own?",
                "choices": {
                    "A": "A sports equipment store",
                    "B": "A winter apparel store",
                    "C": "An automobile dealership",
                    "D": "A hotel chain"
                },
                "answer": "A"
            },
            {
                "id": 58,
                "question": "<strong>Question 58.</strong> According to the men, what has changed recently?",
                "choices": {
                    "A": "Roads have become more accessible.",
                    "B": "Costs have decreased.",
                    "C": "Tourism has increased.",
                    "D": "Weather patterns have shifted."
                },
                "answer": "C"
            }
        ]
    },
    {
        "q_num": "59-61",
        "audio": "E26-T01-59-61.mp3",
        "transcript": [
            "M: Many of our factory workers have expressed interest in upgrading their skills. (59) I'd like to implement a peer-training program, where learners shadow more-experienced employees and observe how they do their jobs.",
            "W: I'm afraid that might become a burden for our long-time employees. (60) They'll have to slow down their work to explain what they're doing.",
            "M: (61) What if we videotaped experienced employees doing specific tasks? High-quality video can be recorded and edited with a smartphone.",
            "W: I like that idea. It would allow us to capture our workers' expertise without slowing down the production line."
        ],
        "questions": [
            {
                "id": 59,
                "question": "<strong>PRACTICE 10</strong><br><strong>Question 59.</strong> What does the man want to do?",
                "choices": {
                    "A": "Provide training opportunities",
                    "B": "Upgrade machinery",
                    "C": "Hire additional employees",
                    "D": "Reorganize the factory layout"
                },
                "answer": "A"
            },
            {
                "id": 60,
                "question": "<strong>Question 60.</strong> What is the woman concerned about?",
                "choices": {
                    "A": "Increasing expenses",
                    "B": "Introducing errors",
                    "C": "Reducing productivity",
                    "D": "Causing confusion"
                },
                "answer": "C"
            },
            {
                "id": 61,
                "question": "<strong>Question 61.</strong> What does the man mean when he says, \"High-quality video can be recorded and edited with a smartphone\"?",
                "choices": {
                    "A": "A new policy should be established.",
                    "B": "An idea is easy to implement.",
                    "C": "Data security is a concern.",
                    "D": "Some information should be verified."
                },
                "answer": "B"
            }
        ]
    },
    {
        "q_num": "62-64",
        "audio": "E26-T01-62-64.mp3",
        "graphic": "data/graphics/ets26_t01_q62_64.png",
        "transcript": [
            "W: Hi, Suresh. (62) I'm at the airport waiting for my flight. I want to meet with a potential investor while I'm in Chicago. Her name's Marta Gomez. I can send you her contact information.",
            "M: OK. (63) Which day would you prefer to meet with her?",
            "W: (63) How about right after my meeting with the Chicago staff?",
            "M: OK. By the way, (64) did you see that our company won an award for our contributions to the community? It was just announced this morning."
        ],
        "questions": [
            {
                "id": 62,
                "question": "<strong>PRACTICE 11 (GRAPHIC)</strong><br><strong>Question 62.</strong> Where is the woman?",
                "choices": {
                    "A": "At a restaurant",
                    "B": "At a travel agency",
                    "C": "At an airport",
                    "D": "At a warehouse"
                },
                "answer": "C"
            },
            {
                "id": 63,
                "question": "<strong>Question 63.</strong> Look at the graphic. When does the woman prefer to meet with an investor?",
                "choices": {
                    "A": "On Tuesday at Noon",
                    "B": "On Wednesday at 8:00 A.M.",
                    "C": "On Wednesday at 2:00 P.M.",
                    "D": "On Thursday at 2:00 P.M."
                },
                "answer": "C"
            },
            {
                "id": 64,
                "question": "<strong>Question 64.</strong> What good news does the man share?",
                "choices": {
                    "A": "A colleague received a promotion.",
                    "B": "A conference proposal was accepted.",
                    "C": "An airline ticket has been upgraded.",
                    "D": "A company won an award."
                },
                "answer": "D"
            }
        ]
    },
    {
        "q_num": "65-67",
        "audio": "E26-T01-65-67.mp3",
        "graphic": "data/graphics/ets26_t01_q65_67.png",
        "transcript": [
            "M: Marion, (65) we keep getting calls from people who want to visit the botanical garden but can't find parking information. Isn't it on our Web site?",
            "W: It is, but you have to click on the \"About Us\" page and scroll to the bottom of that page. Maybe people don't see it.",
            "M: Oh, (66) I think we should move that information from the \"About Us\" page and make a separate page for directions and parking information. That way, people can find it more easily.",
            "W: I'd be happy to make that change. But (67) we're in the middle of updating our software, so it'll have to wait until Monday."
        ],
        "questions": [
            {
                "id": 65,
                "question": "<strong>PRACTICE 12 (GRAPHIC)</strong><br><strong>Question 65.</strong> Where do the speakers work?",
                "choices": {
                    "A": "At an amusement park",
                    "B": "At an art museum",
                    "C": "At a concert hall",
                    "D": "At a botanical garden"
                },
                "answer": "D"
            },
            {
                "id": 66,
                "question": "<strong>Question 66.</strong> Look at the graphic. Which page on the Web site does the man want to change?",
                "choices": {
                    "A": "Page 1",
                    "B": "Page 2",
                    "C": "Page 3",
                    "D": "Page 4"
                },
                "answer": "A"
            },
            {
                "id": 67,
                "question": "<strong>Question 67.</strong> Why does the woman say she cannot complete a task until Monday?",
                "choices": {
                    "A": "She requires approval from a manager.",
                    "B": "She is attending a workshop.",
                    "C": "Some software is being updated.",
                    "D": "Some clients will be arriving soon."
                },
                "answer": "C"
            }
        ]
    },
    {
        "q_num": "68-70",
        "audio": "E26-T01-68-70.mp3",
        "graphic": "data/graphics/ets26_t01_q68_70.png",
        "transcript": [
            "M: Good news! (68) We have finally received the go-ahead for our department's project to install bicycle racks at the train station downtown.",
            "W: At last! So, now we need to decide where to place the racks. How about by the station entrance?",
            "M: Hmm. If we asked riders, (69) I bet they'd say that the most convenient spot is as close to the platform as possible.",
            "W: (69) Let's do that. (70) I'll contact some companies for estimates."
        ],
        "questions": [
            {
                "id": 68,
                "question": "<strong>PRACTICE 13 (GRAPHIC)</strong><br><strong>Question 68.</strong> What news does the man share?",
                "choices": {
                    "A": "A station road will be closed for repair.",
                    "B": "A project has been approved.",
                    "C": "A parking area has been expanded.",
                    "D": "An office will relocate."
                },
                "answer": "B"
            },
            {
                "id": 69,
                "question": "<strong>Question 69.</strong> Look at the graphic. Where do the speakers decide to install some bicycle racks?",
                "choices": {
                    "A": "Near the covered parking area",
                    "B": "Near the long-term parking area",
                    "C": "Near the short-term parking area",
                    "D": "Near the overflow parking area"
                },
                "answer": "A"
            },
            {
                "id": 70,
                "question": "<strong>Question 70.</strong> Why does the woman say she will contact some companies?",
                "choices": {
                    "A": "To arrange a loan",
                    "B": "To apply for a permit",
                    "C": "To ask for estimates",
                    "D": "To create a proposal"
                },
                "answer": "C"
            }
        ]
    }
]

# --- BEGIN ETS 2026 TEST 1 GENERATION ---
from translate_database import translate_text, generate_explanation

# Remove old topic sections
final_data = [s for s in final_data if not s.get("id", "").startswith("topic_")]

# Translate & build practice sets
processed_practice_sets = []
for s_idx, r_set in enumerate(raw_sets):
    vi_transcript = [translate_text(line) for line in r_set["transcript"]]
    processed_questions = []
    for q in r_set["questions"]:
        # Clean up the question text (remove PRACTICE and Question headers)
        q_clean = q["question"]
        q_clean = re.sub(r"<strong>PRACTICE\s*\d+</strong><br>", "", q_clean, flags=re.IGNORECASE)
        q_clean = re.sub(r"<strong>Question\s*\d+[\.:]</strong>\s*", "", q_clean, flags=re.IGNORECASE)
        q_clean = q_clean.strip()
        q["question"] = q_clean
        
        vi_q_text = translate_text(q["question"])
        vi_choices = {}
        for letter, choice_val in q["choices"].items():
            vi_choices[letter] = translate_text(choice_val)
        
        ans_letter = q["answer"]
        key_tag = f"({q['id']})"
        en_key_line = ""
        vi_key_line = ""
        for idx, line in enumerate(r_set["transcript"]):
            if key_tag in line:
                en_key_line = line
                if idx < len(vi_transcript):
                    vi_key_line = vi_transcript[idx]
                break
        if not en_key_line:
            en_key_line = r_set["transcript"][0]
            vi_key_line = vi_transcript[0]
            
        explanation = generate_explanation(q["question"], q["choices"], ans_letter, r_set["transcript"], vi_choices, vi_transcript)
        processed_questions.append({
            "id": q["id"],
            "slide_index": 3000 + q["id"],
            "question": q["question"],
            "choices": q["choices"],
            "answer": q["answer"],
            "vietnamese_question": vi_q_text,
            "vietnamese_choices": vi_choices,
            "explanation": explanation
        })
        
    p_set = {
        "set_index": s_idx + 1,
        "audio": r_set["audio"],
        "questions": processed_questions,
        "transcript": r_set["transcript"],
        "vietnamese_transcript": vi_transcript
    }
    if "graphic" in r_set:
        p_set["graphic"] = r_set["graphic"]
    processed_practice_sets.append(p_set)

new_section = {
    "id": "ets_test_01",
    "title": "ETS 2026 - TEST 01",
    "type": "test",
    "theory": [],
    "vocabulary": [],
    "practice_sets": processed_practice_sets
}
final_data.append(new_section)
# --- END ETS 2026 TEST 1 GENERATION ---

with open(output_json_path, "w", encoding="utf-8") as f:
    json.dump(final_data, f, ensure_ascii=False, indent=2)

# Write JS wrapper file to bypass CORS issues on file:// protocol
output_js_path = os.path.join(data_dir, "part03_data.js")
with open(output_js_path, "w", encoding="utf-8") as f:
    f.write("window.part03Data = ")
    json.dump(final_data, f, ensure_ascii=False, indent=2)
    f.write(";")

# Call the vocabulary loader to re-apply high-quality VSTEP PDF vocabularies
try:
    import load_pdf_vocabulary
    load_pdf_vocabulary.main()
    print("Successfully re-applied VSTEP PDF vocabularies.")
except Exception as e:
    print(f"Warning: Failed to auto-apply VSTEP PDF vocabularies: {e}")

print(f"\nProcessing complete! Compiled data saved to {output_json_path} and {output_js_path}")
print("Ready for web integration!")
