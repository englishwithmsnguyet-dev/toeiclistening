import pptx
prs = pptx.Presentation('../TOEIC LISTENING - PART 01.pptx')
for idx in range(206, 212):
    slide = prs.slides[idx]
    print(f"\nSlide {idx+1}:")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"  Shape {s_idx}: {shape.shape_type} - {shape.name}")
        if hasattr(shape, "image"):
            print(f"    Image size: {len(shape.image.blob)} bytes, ext: {shape.image.ext}")
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for gs in shape.shapes:
                print(f"    Group Shape: {gs.shape_type} - {gs.name}")
                if hasattr(gs, "image"):
                    print(f"      Image size: {len(gs.image.blob)} bytes, ext: {gs.image.ext}")
