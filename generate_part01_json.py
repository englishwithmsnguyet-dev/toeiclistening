import os
import json
import re
import pptx
import fitz

PPTX_PATH = '../TOEIC LISTENING - PART 01.pptx'
PDF_PATH = '/Users/nguyetpham/Desktop/TEACHING/TOEIC 2026/TÀI LIỆU ETS/ETS 2026/TOEIC ETS 2026- Thay Dinh Van/Final Thay Dinh Van LISTENING ETS 2026 .pdf'
OUTPUT_JSON = 'data/part01_data.json'
OUTPUT_JS = 'data/part01_data.js'
GRAPHICS_DIR = 'data/graphics/part01'

if not os.path.exists(GRAPHICS_DIR):
    os.makedirs(GRAPHICS_DIR)

def get_paragraph_html(p):
    text_html = ""
    for run in p.runs:
        t = run.text.replace('<', '&lt;').replace('>', '&gt;')
        if not t: continue
        if run.font.bold:
            t = f"<strong>{t}</strong>"
        if run.font.italic:
            t = f"<em>{t}</em>"
        
        if run.font.color and run.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
            hex_color = str(run.font.color.rgb)
            if hex_color and hex_color != "000000":
                t = f"<span style=\"color: #{hex_color};\">{t}</span>"
        
        text_html += t
    return text_html

def extract_slide_text(slide):
    sorted_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: s.top
    )
    texts = []
    for shape in sorted_shapes:
        for p in shape.text_frame.paragraphs:
            p_html = get_paragraph_html(p)
            if p_html.strip():
                is_bullet = False
                pPr = p._p.pPr
                if pPr is not None:
                    if (pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar') is not None or 
                        pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buFont') is not None) and \
                       pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone') is None:
                        is_bullet = True
                
                text_plain = p.text.strip()
                if not is_bullet and text_plain.startswith(('o ', '• ', '- ', '* ', '◦ ')):
                    is_bullet = True
                
                if is_bullet and not text_plain.startswith(('o ', '• ', '- ', '* ', '◦ ')):
                    p_html = "• " + p_html
                        
                texts.append(p_html)
    return texts

def extract_slide_audio(slide):
    for shape in slide.shapes:
        xml_str = shape.element.xml
        if 'media' in xml_str or 'audio' in xml_str:
            rids = re.findall(r'r:(?:embed|id|link)="([^"]+)"', xml_str)
            for rid in rids:
                try:
                    target = slide.part.rels[rid].target_ref
                    if target.endswith('.mp3'):
                        return target.split('/')[-1]
                except:
                    pass
    return None

def extract_slide_images(slide, slide_num):
    imgs = []
    for idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "image") and len(shape.image.blob) > 20000: # ignore tiny icons
            ext = shape.image.ext
            filename = f"slide_{slide_num}_img_{idx}.{ext}"
            filepath = os.path.join(GRAPHICS_DIR, filename)
            with open(filepath, "wb") as f:
                f.write(shape.image.blob)
            imgs.append(filename)
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for gidx, gs in enumerate(shape.shapes):
                if hasattr(gs, "image") and len(gs.image.blob) > 20000:
                    ext = gs.image.ext
                    filename = f"slide_{slide_num}_img_{idx}_{gidx}.{ext}"
                    filepath = os.path.join(GRAPHICS_DIR, filename)
                    with open(filepath, "wb") as f:
                        f.write(gs.image.blob)
                    imgs.append(filename)
    return imgs

def parse_theory_slides(prs, slide_indices):
    data = []
    for s_idx in slide_indices:
        slide = prs.slides[s_idx - 1]
        data.append({
            "slide_index": s_idx,
            "text": extract_slide_text(slide),
            "images": extract_slide_images(slide, s_idx),
            "audio": extract_slide_audio(slide)
        })
    return data

def find_test_answer(slide):
    # Find the oval shape and its top coordinate
    for s in slide.shapes:
        if 'Oval' in s.name:
            top = s.top
            if top < 4500000: return 'A'
            elif top < 6000000: return 'B'
            elif top < 7000000: return 'C'
            else: return 'D'
    return 'A' # fallback

# Structure map
structure = [
    {"id": "overview", "title": "I. TỔNG QUAN PHẦN 01", "type": "overview", "slides": list(range(4, 10))},
    {"id": "dang_01", "title": "DẠNG 01: CÓ MỘT NGƯỜI TRONG HÌNH", "type": "theory", "slides": list(range(10, 111))},
    {"id": "dang_02", "title": "DẠNG 02: TRANH CÓ NHIỀU NGƯỜI", "type": "theory", "slides": list(range(111, 149))},
    {"id": "dang_03", "title": "DẠNG 03: TRANH MIÊU TẢ VẬT", "type": "theory", "slides": list(range(149, 206))}
]

print("Loading PPTX...")
prs = pptx.Presentation(PPTX_PATH)
final_data = []

for sec in structure:
    print(f"Parsing {sec['title']}...")
    final_data.append({
        "id": sec["id"],
        "title": sec["title"],
        "type": sec["type"],
        "theory": parse_theory_slides(prs, sec["slides"])
    })

# Extract test photos from PDF
print("Loading PDF for Test Photographs...")
doc = fitz.open(PDF_PATH)
def extract_test_photos(test_idx):
    # Test 1 starts at page 22 (idx 21), Test 2 at page 32 (idx 31)...
    start_p = 21 + (test_idx - 1) * 10
    photos = []
    for p in range(start_p, start_p + 3):
        imgs = doc[p].get_images()
        big = [(img[0], img[2], img[3]) for img in imgs if img[2] > 2000 and img[3] > 1000]
        # Sort by vertical position or assume they are naturally ordered. PyMuPDF usually returns them in draw order.
        for xref, w, h in big:
            pix = fitz.Pixmap(doc, xref)
            filename = f"ets26_t{test_idx:02d}_q{len(photos)+1:02d}.jpg"
            filepath = os.path.join(GRAPHICS_DIR, filename)
            pix.save(filepath)
            photos.append(filename)
            if len(photos) == 6: break
    # Note: On page 1 of each test, the first big image is the Example Photo. We must skip it!
    # The actual 6 photos are photos[1:7] if we extracted 7.
    # Wait, earlier we saw page 1 had 3 images (img0=Example, img1=Q1, img2=Q2). 
    # Let's refine extraction:
    actual_photos = []
    # Page 1:
    imgs_p1 = [img for img in doc[start_p].get_images() if img[2] > 2000 and img[3] > 1000]
    for i, img in enumerate(imgs_p1):
        if i == 0: continue # Skip example photo
        pix = fitz.Pixmap(doc, img[0])
        fname = f"ets26_t{test_idx:02d}_q{len(actual_photos)+1:02d}.jpg"
        pix.save(os.path.join(GRAPHICS_DIR, fname))
        actual_photos.append(fname)
    # Page 2:
    imgs_p2 = [img for img in doc[start_p+1].get_images() if img[2] > 2000 and img[3] > 1000]
    for img in imgs_p2:
        pix = fitz.Pixmap(doc, img[0])
        fname = f"ets26_t{test_idx:02d}_q{len(actual_photos)+1:02d}.jpg"
        pix.save(os.path.join(GRAPHICS_DIR, fname))
        actual_photos.append(fname)
    # Page 3:
    imgs_p3 = [img for img in doc[start_p+2].get_images() if img[2] > 2000 and img[3] > 1000]
    for img in imgs_p3:
        pix = fitz.Pixmap(doc, img[0])
        fname = f"ets26_t{test_idx:02d}_q{len(actual_photos)+1:02d}.jpg"
        pix.save(os.path.join(GRAPHICS_DIR, fname))
        actual_photos.append(fname)
        if len(actual_photos) == 6: break
    
    return actual_photos

print("Parsing Tests 1-5...")
for t_idx in range(1, 6):
    print(f"  Test {t_idx}...")
    photos = extract_test_photos(t_idx)
    
    # Slides for this test: 
    # Test 1: 207-212
    # Test 2: 214-219
    # Test 3: 221-226
    # Test 4: 228-233
    # Test 5: 235-240
    start_s = 206 + (t_idx - 1) * 7
    
    practice_sets = []
    for q_idx in range(6):
        s_idx = start_s + q_idx
        slide = prs.slides[s_idx]
        
        # Extract statements
        # The text shape is usually the one with the most text
        txt_shape = max([s for s in slide.shapes if s.has_text_frame], key=lambda s: len(s.text_frame.text))
        paras = [p.text.strip() for p in txt_shape.text_frame.paragraphs if p.text.strip()]
        
        # Typically paras are: "1. Statement A", "(B) Statement B", "(C) ...", "(D) ..."
        choices = {"A": "", "B": "", "C": "", "D": ""}
        for p in paras:
            p_clean = re.sub(r'^(\d+\.|[\(]?[A-D][\)]?)\s*', '', p).strip()
            if '(B)' in p or p.startswith('B.') or (len(paras)==4 and p == paras[1]): choices['B'] = p_clean
            elif '(C)' in p or p.startswith('C.') or (len(paras)==4 and p == paras[2]): choices['C'] = p_clean
            elif '(D)' in p or p.startswith('D.') or (len(paras)==4 and p == paras[3]): choices['D'] = p_clean
            else: choices['A'] = p_clean

        ans = find_test_answer(slide)
        audio_file = f"E26-T{t_idx:02d}-0{q_idx+1}.mp3"
        
        practice_sets.append({
            "set_index": q_idx + 1,
            "audio": audio_file,
            "image": f"part01/{photos[q_idx]}",
            "questions": [{
                "id": q_idx + 1,
                "slide_index": s_idx + 1,
                "question": "Look at the picture and choose the statement that best describes it.",
                "choices": choices,
                "answer": ans,
                "explanation": f"<strong style='color: var(--success);'>ĐÁP ÁN ĐÚNG LÀ {ans}</strong>"
            }]
        })
    
    final_data.append({
        "id": f"test_0{t_idx}",
        "title": f"TEST {t_idx}",
        "type": "test",
        "practice_sets": practice_sets
    })

print(f"Writing to {OUTPUT_JSON} and {OUTPUT_JS}...")
with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
    json.dump(final_data, f, ensure_ascii=False, indent=2)

with open(OUTPUT_JS, "w", encoding="utf-8") as f:
    f.write("window.part01Data = ")
    json.dump(final_data, f, ensure_ascii=False, indent=2)
    f.write(";\n")

print("Done!")
