import sys
import os
import re
import json
import zipfile
import html
from pptx import Presentation

# 1. Paths configuration
script_dir = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.abspath(os.path.join(script_dir, "..", "TOEIC PART 04", "TOEIC PART 04 - 2023.pptx"))
output_dir = script_dir
data_dir = os.path.join(output_dir, "data")
media_dir = os.path.join(output_dir, "media")

# Create folders if not exist
os.makedirs(data_dir, exist_ok=True)
os.makedirs(media_dir, exist_ok=True)

print("Starting TOEIC Listening Part 4 HTML Compiler...")

# 2. Extract media from PPTX zip structure
with zipfile.ZipFile(pptx_path, 'r') as z:
    audio_files = [f for f in z.namelist() if 'media' in f and f.lower().endswith('.mp3')]
    print(f"Found {len(audio_files)} audio files in PPTX.")
    extracted_count = 0
    for f in audio_files:
        filename = os.path.basename(f)
        dest_path = os.path.join(media_dir, filename)
        if not os.path.exists(dest_path):
            with open(dest_path, 'wb') as dest:
                dest.write(z.read(f))
            extracted_count += 1
    print(f"Extracted {extracted_count} audio files to {media_dir}.")

# Load presentation
prs = Presentation(pptx_path)
slides = list(prs.slides)
print(f"Loaded presentation with {len(slides)} slides.")

# Helpers
def get_paragraph_html(paragraph):
    parts = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        
        t_escaped = html.escape(text)
        color = run.font.color
        has_color = False
        hex_color = ""
        if color and color.type == 1:
            try:
                hex_color = f"#{color.rgb}"
                has_color = True
            except:
                pass
                
        span_text = t_escaped
        style_attrs = []
        if has_color:
            style_attrs.append(f"color: {hex_color};")
            
        if style_attrs:
            style_str = " ".join(style_attrs)
            span_text = f'<span style="{style_str}">{span_text}</span>'
            
        if run.font.bold:
            span_text = f'<strong>{span_text}</strong>'
        if run.font.italic:
            span_text = f'<em>{span_text}</em>'
            
        span_text = span_text.replace('\x0b', '<br>')
        span_text = span_text.replace('\u000b', '<br>')
        parts.append(span_text)
        
    return "".join(parts).strip()

def extract_slide_text(slide):
    sorted_shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: s.top
    )
    texts = []
    for shape in sorted_shapes:
        for p in shape.text_frame.paragraphs:
            p_html = get_paragraph_html(p)
            if p_html:
                is_bullet = False
                pPr = p._p.pPr
                if pPr is not None:
                    has_bu_char = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar') is not None
                    has_bu_font = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buFont') is not None
                    has_bu_none = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone') is not None
                    if (has_bu_char or has_bu_font) and not has_bu_none:
                        is_bullet = True
                        
                text_plain = p.text.strip()
                if not is_bullet:
                    if text_plain.startswith(('o ', '• ', '- ', '* ', '◦ ')):
                        is_bullet = True
                
                if is_bullet:
                    if not text_plain.startswith(('o ', '• ', '- ', '* ', '◦ ')):
                        p_html = "• " + p_html
                texts.append(p_html)
    return texts

def extract_slide_audio(slide):
    for shape in slide.shapes:
        xml_str = shape.element.xml
        if 'media' in xml_str or 'audio' in xml_str:
            rids = re.findall(r'r:(?:embed|id|link)="([^"]+)"', xml_str)
            for rid in rids:
                try:
                    target = slide.part.rels[rid].target_ref
                    if target.endswith('.mp3'):
                        return target.split('/')[-1]
                except:
                    pass
    return None

def find_nearest_audio(slides, target_slide_num):
    start = max(1, target_slide_num - 5)
    end = min(len(slides), target_slide_num + 5)
    for s_idx in range(start, end + 1):
        audio = extract_slide_audio(slides[s_idx - 1])
        if audio:
            return audio
    return None

def parse_question_slide(slide, slide_num):
    tb = None
    for shape in slide.shapes:
        if shape.has_text_frame and "A." in shape.text_frame.text:
            tb = shape
            break
            
    if not tb:
        return None
        
    p_htmls = []
    for p in tb.text_frame.paragraphs:
        p_html = get_paragraph_html(p)
        if p_html:
            p_htmls.append(p_html)
            
    v_lines_with_col = []
    for p_html in p_htmls:
        lines = re.split(r'<br\s*/?>', p_html)
        for line in lines:
            subparts = re.split(r'\t|\s{2,}(?=[A-D]\.)|\s+(?=[B-D]\.)', line)
            subparts = [s.strip() for s in subparts if s.strip()]
            for col_idx, part in enumerate(subparts):
                v_lines_with_col.append((part, col_idx))
                
    question_text = ""
    choices = {}
    choice_lines = []
    
    for l_idx, (line_html, col_idx) in enumerate(v_lines_with_col):
        text_plain = re.sub(r'<[^>]+>', '', line_html).strip()
        match = re.match(r'^([A-D])\.\s*(.*)', text_plain)
        if match:
            letter = match.group(1)
            html_clean = re.sub(r'^((?:<[^>]+>)*)([A-D])\.\s*', r'\1', line_html)
            choices[letter] = html_clean
            choice_lines.append((l_idx, letter, col_idx))
        else:
            if not question_text:
                question_text = line_html
            else:
                question_text += "<br>" + line_html
                
    oval_shape = None
    for shape in slide.shapes:
        if "oval" in shape.name.lower() or "circle" in shape.name.lower():
            oval_shape = shape
            break
            
    correct_answer = None
    if oval_shape and choice_lines:
        tb_top = tb.top
        tb_height = tb.height
        tb_left = tb.left
        tb_width = tb.width
        oval_top = oval_shape.top
        oval_left = oval_shape.left
        
        oval_x = oval_left + oval_shape.width / 2
        is_left_side = oval_x < (tb_left + tb_width * 0.48)
        target_col = 0 if is_left_side else 1
        has_target_col_choices = any(col == target_col for _, _, col in choice_lines)
        
        best_letter = None
        min_dist = float('inf')
        for l_idx, letter, col_idx in choice_lines:
            if has_target_col_choices:
                if col_idx != target_col:
                    continue
            y_est = tb_top + ((l_idx + 0.5) / len(v_lines_with_col)) * tb_height
            dist = abs(oval_top - y_est)
            if dist < min_dist:
                min_dist = dist
                best_letter = letter
        correct_answer = best_letter
        
    audio = extract_slide_audio(slide)
    
    return {
        "slide_index": slide_num,
        "question": question_text,
        "choices": choices,
        "answer": correct_answer,
        "audio": audio
    }

def get_slide_script_html(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            p_htmls = []
            for p in shape.text_frame.paragraphs:
                p_html = get_paragraph_html(p)
                if p_html:
                    p_htmls.append(p_html)
            return p_htmls
    return []

# Structure definition for Part 04 (1-based slide indices)
structure_map = [
    {
        "id": "overview",
        "title": "TỔNG QUAN PHẦN 04",
        "type": "overview",
        "slides": [2, 3]
    },
    {
        "id": "topic_01",
        "title": "Chủ đề 1 - Tin nhắn điện thoại",
        "type": "topic",
        "theory_slides": [4, 5, 6, 7, 8],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [9, 10, 11, 12]},
            {"slides": [13, 14, 15, 16]}
        ],
        "practice_sets": [
            {"slides": list(range(18, 25))},
            {"slides": list(range(25, 32))},
            {"slides": list(range(32, 39))},
            {"slides": list(range(39, 46))},
            {"slides": list(range(46, 53))},
            {"slides": list(range(53, 60))}
        ]
    },
    {
        "id": "topic_02",
        "title": "Chủ đề 2 - Thông báo",
        "type": "topic",
        "theory_slides": [60, 61, 62],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [63, 64, 65, 66]},
            {"slides": [67, 68, 69, 70]}
        ],
        "practice_sets": [
            {"slides": list(range(72, 79))},
            {"slides": list(range(79, 86))},
            {"slides": list(range(86, 93))},
            {"slides": list(range(93, 100))},
            {"slides": list(range(100, 107))},
            {"slides": list(range(107, 114))}
        ]
    },
    {
        "id": "topic_03",
        "title": "Chủ đề 3 - Quảng cáo",
        "type": "topic",
        "theory_slides": [114, 115, 116, 117],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [118, 119, 120, 121]},
            {"slides": [122, 123, 124, 125]}
        ],
        "practice_sets": [
            {"slides": list(range(127, 134))},
            {"slides": list(range(134, 141))},
            {"slides": list(range(141, 148))},
            {"slides": list(range(148, 155))},
            {"slides": list(range(155, 162))},
            {"slides": list(range(162, 169))}
        ]
    },
    {
        "id": "topic_04",
        "title": "Chủ đề 4 - Buổi phát thanh trên radio",
        "type": "topic",
        "theory_slides": [169, 170, 171],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [172, 173, 174, 175]},
            {"slides": [176, 177, 178, 179]}
        ],
        "practice_sets": [
            {"slides": list(range(181, 188))},
            {"slides": list(range(188, 195))},
            {"slides": list(range(195, 202))},
            {"slides": list(range(202, 209))},
            {"slides": list(range(209, 216))}
        ]
    },
    {
        "id": "topic_05",
        "title": "Chủ đề 5 - Bài diễn thuyết và phát biểu",
        "type": "topic",
        "theory_slides": [216, 217, 218, 219, 220],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [221, 222, 223, 224]},
            {"slides": [225, 226, 227, 228]}
        ],
        "practice_sets": [
            {"slides": list(range(230, 237))},
            {"slides": list(range(237, 244))},
            {"slides": list(range(244, 251))},
            {"slides": list(range(251, 258))},
            {"slides": list(range(258, 265))},
            {"slides": list(range(265, 272))}
        ]
    },
    {
        "id": "topic_06",
        "title": "Chủ đề 6 - Trích đoạn từ cuộc họp",
        "type": "topic",
        "theory_slides": [272, 273],
        "vocabulary_slides": [],
        "example_sets": [
            {"slides": [274, 275, 276, 277]},
            {"slides": [278, 279, 280, 281]}
        ],
        "practice_sets": [
            {"slides": list(range(283, 290))},
            {"slides": list(range(290, 297))},
            {"slides": list(range(297, 304))},
            {"slides": list(range(304, 311))},
            {"slides": list(range(311, 318))},
            {"slides": list(range(318, 325))}
        ]
    }
]

# 4. Processing Slide Content
final_data = []

for entry in structure_map:
    print(f"Processing Section: {entry['title']}...")
    
    if entry["type"] == "overview" or entry["type"] == "tips":
        theory_data = []
        for s_num in entry["slides"]:
            slide = slides[s_num - 1]
            theory_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
        final_data.append({
            "id": entry["id"],
            "title": entry["title"],
            "type": entry["type"],
            "theory": theory_data,
            "vocabulary": []
        })
        
    elif entry["type"] == "topic":
        theory_data = []
        for s_num in entry["theory_slides"]:
            slide = slides[s_num - 1]
            theory_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
            
        vocabulary_data = []
        for s_num in entry["vocabulary_slides"]:
            slide = slides[s_num - 1]
            vocabulary_data.append({
                "slide_index": s_num,
                "text": extract_slide_text(slide)
            })
            
        example_sets_data = []
        for set_idx, ex_set in enumerate(entry.get("example_sets", [])):
            set_slides_nums = ex_set["slides"]
            
            q1_num = set_slides_nums[0]
            q2_num = set_slides_nums[1]
            q3_num = set_slides_nums[2]
            full_s_num = set_slides_nums[3]
            
            q1_info = parse_question_slide(slides[q1_num - 1], q1_num)
            q2_info = parse_question_slide(slides[q2_num - 1], q2_num)
            q3_info = parse_question_slide(slides[q3_num - 1], q3_num)
            
            full_transcript = extract_slide_text(slides[full_s_num - 1])
            
            audio = None
            for s_num in set_slides_nums:
                slide_audio = extract_slide_audio(slides[s_num - 1])
                if slide_audio:
                    audio = slide_audio
                    break
                    
            questions = []
            for i, qi in enumerate([q1_info, q2_info, q3_info]):
                if qi:
                    questions.append({
                        "id": i + 1,
                        "slide_index": qi["slide_index"],
                        "question": qi["question"],
                        "choices": qi["choices"],
                        "answer": qi["answer"]
                    })
                    
            example_sets_data.append({
                "set_index": set_idx + 1,
                "audio": audio,
                "questions": questions,
                "transcript": full_transcript
            })

        practice_sets_data = []
        for set_idx, pr_set in enumerate(entry["practice_sets"]):
            set_slides_nums = pr_set["slides"]
            
            q1_num = set_slides_nums[0]
            s1_num = set_slides_nums[1]
            q2_num = set_slides_nums[2]
            s2_num = set_slides_nums[3]
            q3_num = set_slides_nums[4]
            s3_num = set_slides_nums[5]
            full_s_num = set_slides_nums[6]
            
            q1_info = parse_question_slide(slides[q1_num - 1], q1_num)
            q2_info = parse_question_slide(slides[q2_num - 1], q2_num)
            q3_info = parse_question_slide(slides[q3_num - 1], q3_num)
            
            full_transcript = extract_slide_text(slides[full_s_num - 1])
            
            audio = None
            for s_num in set_slides_nums:
                slide_audio = extract_slide_audio(slides[s_num - 1])
                if slide_audio:
                    audio = slide_audio
                    break
                
            questions = []
            for i, qi in enumerate([q1_info, q2_info, q3_info]):
                if qi:
                    questions.append({
                        "id": i + 1,
                        "slide_index": qi["slide_index"],
                        "question": qi["question"],
                        "choices": qi["choices"],
                        "answer": qi["answer"]
                    })
                    
            practice_sets_data.append({
                "set_index": set_idx + 1,
                "audio": audio,
                "questions": questions,
                "transcript": full_transcript
            })
            
        final_data.append({
            "id": entry["id"],
            "title": entry["title"],
            "type": "topic",
            "theory": theory_data,
            "vocabulary": vocabulary_data,
            "examples": example_sets_data,
            "practice_sets": practice_sets_data
        })

# 5. Output file writing
output_json_path = os.path.join(data_dir, "part04_data.json")
with open(output_json_path, "w", encoding="utf-8") as f:
    json.dump(final_data, f, ensure_ascii=False, indent=4)
print(f"Data successfully saved to {output_json_path}")

output_js_path = os.path.join(data_dir, "part04_data.js")
js_content = f"window.part04Data = {json.dumps(final_data, ensure_ascii=False, indent=2)};\n"
with open(output_js_path, "w", encoding="utf-8") as f:
    f.write(js_content)
print(f"JS Data successfully saved to {output_js_path}")
print("Done!")
