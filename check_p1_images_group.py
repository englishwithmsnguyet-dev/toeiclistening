import pptx
prs = pptx.Presentation('../TOEIC LISTENING - PART 01.pptx')
for idx in range(206, 212):
    slide = prs.slides[idx]
    has_img = False
    for shape in slide.shapes:
        if hasattr(shape, "image") and len(shape.image.blob) > 20000:
            has_img = True
            print(f"Slide {idx+1} has image: {len(shape.image.blob)} bytes")
        elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for gs in shape.shapes:
                if hasattr(gs, "image") and len(gs.image.blob) > 20000:
                    has_img = True
                    print(f"Slide {idx+1} has grouped image: {len(gs.image.blob)} bytes")
    if not has_img:
        print(f"Slide {idx+1} has NO image!")
